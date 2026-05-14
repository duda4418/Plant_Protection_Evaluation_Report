from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def generate_approach_slides(output_path: Path) -> None:
    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation = Presentation()
    _set_wide_layout(presentation)

    slides = [
        (
            "Dynamic Regulatory Report Generator",
            [
                "Python prototype for ActiveDocs-style document automation",
                "JSON input validated, mapped, visualized, and rendered into DOCX",
                "Built for maintainable report templates and explicit business rules",
            ],
        ),
        (
            "End-to-End Flow",
            [
                "input_data.json -> Pydantic validation -> view-model mapping",
                "Business rules stay in Python and are reusable outside the template",
                "Charts are generated before DOCX rendering and embedded into the report",
            ],
        ),
        (
            "Template Architecture",
            [
                "The product section is a reusable repeating template component",
                "Tables use row-level loops for applications, toxicity, and risk components",
                "Headers, footers, classification, and page fields are part of the template",
            ],
        ),
        (
            "Business Logic",
            [
                "High-risk notice appears only when risk_score > 70",
                "Restrictions appear only for conditional approvals",
                "German country wording is applied when country == DE",
            ],
        ),
        (
            "Quality Controls",
            [
                "Strict schema validation catches missing fields and invalid ranges",
                "StrictUndefined fails rendering when a template variable is missing",
                "Generated DOCX is checked for unresolved Jinja placeholders",
            ],
        ),
        (
            "Scaling Pattern",
            [
                "Add products or applications without changing the template structure",
                "Keep raw integration data separate from report-ready display context",
                "Large datasets can be summarized in DOCX and exported separately as detail files",
            ],
        ),
        (
            "ActiveDocs Translation",
            [
                "Pydantic models become the data contract or input schema",
                "Mapper functions become the data preparation layer",
                "Jinja loops and conditions map to ActiveDocs repeating sections and rules",
            ],
        ),
    ]

    for title, bullets in slides:
        _add_bullet_slide(presentation, title, bullets)

    presentation.save(output_path)


def _set_wide_layout(presentation: Presentation) -> None:
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)


def _add_bullet_slide(presentation: Presentation, title: str, bullets: list[str]) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(11.9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.clear()
    paragraph = title_frame.paragraphs[0]
    paragraph.text = title
    paragraph.font.name = "Aptos Display"
    paragraph.font.size = Pt(31)
    paragraph.font.bold = True
    paragraph.alignment = PP_ALIGN.LEFT

    body_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.65), Inches(11.2), Inches(4.8))
    body_frame = body_box.text_frame
    body_frame.word_wrap = True
    body_frame.clear()

    for index, bullet in enumerate(bullets):
        item = body_frame.paragraphs[0] if index == 0 else body_frame.add_paragraph()
        item.text = bullet
        item.level = 0
        item.font.name = "Aptos"
        item.font.size = Pt(22)
        item.space_after = Pt(16)

    footer = slide.shapes.add_textbox(Inches(0.7), Inches(6.85), Inches(11.9), Inches(0.3))
    footer_frame = footer.text_frame
    footer_frame.clear()
    footer_para = footer_frame.paragraphs[0]
    footer_para.text = "Plant Protection Evaluation Report | Document automation prototype"
    footer_para.font.name = "Aptos"
    footer_para.font.size = Pt(9)
